# chat/gdrive_ingest.py — robust Google Drive ingestion with better error handling
import io
import json
import time
from typing import Iterable, Tuple
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.oauth2 import service_account
from google.auth.transport.requests import Request
from google.auth import default, compute_engine
from django.conf import settings
from pdfminer.high_level import extract_text as pdf_extract
from pypdf import PdfReader
from docx import Document
import socket

# Drive scopes: read-only
SCOPES = [
    "https://www.googleapis.com/auth/drive.readonly",
    "https://www.googleapis.com/auth/drive.metadata.readonly"
]

MIME_EXPORTS = {
    "application/vnd.google-apps.document": ("text/plain", "txt"),
    "application/vnd.google-apps.presentation": ("text/plain", "txt"),
    "application/vnd.google-apps.spreadsheet": ("text/plain", "txt"),
}
FOLDER_MIME = "application/vnd.google-apps.folder"
TEXT_MIMES = {
    "text/plain",
    "text/markdown",
}

PDF_MIME = "application/pdf"

# Microsoft Office document MIME types
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

# Legacy Office formats
DOC_MIME = "application/msword"
XLS_MIME = "application/vnd.ms-excel"
PPT_MIME = "application/vnd.ms-powerpoint"


def extract_docx_text(file_data: bytes) -> str:
    """Extract text from a .docx file preserving structure."""
    try:
        doc = Document(io.BytesIO(file_data))
        content = []
        
        for element in doc.element.body:
            # Handle paragraphs
            if element.tag.endswith('p'):
                paragraph = None
                for p in doc.paragraphs:
                    if p._element == element:
                        paragraph = p
                        break
                
                if paragraph and paragraph.text.strip():
                    text = paragraph.text.strip()
                    
                    # Detect and preserve list items
                    if text.startswith('•') or text.startswith('-') or text.startswith('*'):
                        content.append(f"• {text[1:].strip()}")
                    elif any(text.startswith(f"{i}.") for i in range(1, 20)):
                        content.append(text)  # Keep numbered lists as-is
                    else:
                        content.append(text)
            
            # Handle tables
            elif element.tag.endswith('tbl'):
                for table in doc.tables:
                    if table._element == element:
                        table_text = extract_table_from_docx(table)
                        if table_text:
                            content.append(table_text)
                        break
        
        return "\n\n".join(content)
    except Exception as e:
        print(f"Error extracting DOCX text: {e}")
        return ""

def extract_table_from_docx(table) -> str:
    """Extract table data in a structured format."""
    try:
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace('\n', ' ')
                cells.append(cell_text)
            rows.append(" | ".join(cells))
        
        if rows:
            # Add markdown table formatting
            header_separator = " | ".join(["---"] * len(rows[0].split(" | ")))
            if len(rows) > 1:
                return f"{rows[0]}\n{header_separator}\n" + "\n".join(rows[1:])
            else:
                return rows[0]
        return ""
    except Exception as e:
        print(f"Error extracting table: {e}")
        return ""


def extract_xlsx_text(file_data: bytes) -> str:
    """Extract text from an Excel file by reading all sheets and cells."""
    try:
        import openpyxl
        workbook = openpyxl.load_workbook(io.BytesIO(file_data), data_only=True)
        all_text = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = []
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                row_text = [cell for cell in row_text if cell.strip()]
                if row_text:
                    sheet_text.append(" | ".join(row_text))
            
            if sheet_text:
                all_text.append(f"Sheet: {sheet_name}\n" + "\n".join(sheet_text))
        
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"Error extracting XLSX text: {e}")
        return ""


def extract_pptx_text(file_data: bytes) -> str:
    """Extract text from a PowerPoint file."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_data))
        all_text = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                all_text.append(f"Slide {i}:\n" + "\n".join(slide_text))
        
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"Error extracting PPTX text: {e}")
        return ""


def _creds_with_retry_old(max_retries=3):
    """Create credentials using Service Account with retry logic."""
    for attempt in range(max_retries):
        try:
            print(f"Attempting to authenticate... (attempt {attempt + 1}/{max_retries})")
            
            # Set socket timeout to prevent hanging
            socket.setdefaulttimeout(30)
            
            credentials = service_account.Credentials.from_service_account_file(
                settings.GDRIVE_SERVICE_ACCOUNT_PATH,
                scopes=SCOPES
            )
            
            # Force refresh to test connection
            request = Request()
            credentials.refresh(request)
            
            print("✅ Authentication successful!")
            return credentials
            
        except socket.timeout:
            print(f"❌ Connection timeout on attempt {attempt + 1}")
            if attempt < max_retries - 1:
                wait_time = (attempt + 1) * 5
                print(f"Waiting {wait_time} seconds before retry...")
                time.sleep(wait_time)
            else:
                raise Exception("Authentication failed after multiple timeout attempts. Check your internet connection.")
                
        except FileNotFoundError:
            raise Exception(f"Service account file not found at: {settings.GDRIVE_SERVICE_ACCOUNT_PATH}")
            
        except Exception as e:
            if "invalid_grant" in str(e):
                raise Exception("Invalid service account credentials. Please check your service_account.json file.")
            elif "access_denied" in str(e):
                raise Exception("Access denied. Make sure your service account has access to the Drive folder.")
            else:
                print(f"❌ Authentication error on attempt {attempt + 1}: {e}")
                if attempt < max_retries - 1:
                    time.sleep(2)
                else:
                    raise

def _creds_with_retry(max_retries=3):
    """Create credentials using ADC (works locally and in GCP) with retry logic."""
    for attempt in range(max_retries):
        try:
            print(f"Attempting to authenticate... (attempt {attempt + 1}/{max_retries})")

            # Set socket timeout to prevent hanging
            socket.setdefaulttimeout(30)

            # Use Application Default Credentials (ADC)
            if settings.APP_ENV == "local":
                print("Using local ADC credentials (gcloud auth application-default login)")
                credentials, _ = default(scopes=SCOPES)
            else:
                print("Using Compute Engine credentials (service account attached to VM/container)")
                credentials = compute_engine.Credentials()

            # Force refresh to test connection
            request = Request()
            credentials.refresh(request)

            print("✅ Authentication successful!")
            return credentials

        except socket.timeout:
            print(f"❌ Connection timeout on attempt {attempt + 1}")
            if attempt < max_retries - 1:
                wait_time = (attempt + 1) * 5
                print(f"Waiting {wait_time} seconds before retry...")
                time.sleep(wait_time)
            else:
                raise Exception("Authentication failed after multiple timeout attempts. Check your internet connection.")

        except Exception as e:
            if "invalid_grant" in str(e):
                raise Exception("Invalid credentials. If local, run `gcloud auth application-default login`. If in GCP, make sure a service account is attached.")
            elif "access_denied" in str(e):
                raise Exception("Access denied. Make sure your identity (user or service account) has access to the Drive folder.")
            else:
                print(f"❌ Authentication error on attempt {attempt + 1}: {e}")
                if attempt < max_retries - 1:
                    time.sleep(2)
                else:
                    raise


def iter_drive_files(folder_id: str):
    """
    Recursively iterate all files in a Google Drive folder,
    including files inside nested folders.
    """

    credentials = _creds_with_retry()
    service = build("drive", "v3", credentials=credentials)

    visited_folders = set()

    def walk(current_folder_id):
        if current_folder_id in visited_folders:
            return

        visited_folders.add(current_folder_id)

        page_token = None

        while True:
            resp = service.files().list(
                q=f"'{current_folder_id}' in parents and trashed=false",
                fields="nextPageToken, files(id,name,mimeType)",
                pageToken=page_token,
                supportsAllDrives=True,
                includeItemsFromAllDrives=True,
            ).execute()

            for item in resp.get("files", []):

                # Recurse into subfolders
                if item["mimeType"] == FOLDER_MIME:
                    print(f"📁 Entering folder: {item['name']}")
                    yield from walk(item["id"])
                else:
                    yield item

            page_token = resp.get("nextPageToken")

            if not page_token:
                break

    # Verify root folder access
    folder_info = service.files().get(
        fileId=folder_id,
        fields="id,name",
        supportsAllDrives=True,
    ).execute()

    print(f"✅ Successfully accessed folder: {folder_info['name']}")

    yield from walk(folder_id)


def fetch_text_for_file(file: dict) -> Tuple[str, dict]:
    """Returns (text, metadata) for a given Drive file with retry logic."""
    try:
        credentials = _creds_with_retry()
        service = build("drive", "v3", credentials=credentials)
        
        fid = file["id"]
        name = file["name"]
        mime = file["mimeType"]

        meta = {"id": fid, "name": name, "mime": mime}

        # Google Docs/Sheets/Slides: export as text
        if mime in MIME_EXPORTS:
            export_mime, ext = MIME_EXPORTS[mime]
            try:
                data = service.files().export(fileId=fid, mimeType=export_mime).execute()
                text = data.decode("utf-8", errors="ignore")
                return text, meta
            except Exception as e:
                print(f"Error exporting {name}: {e}")
                return "", meta

        # Plain text / Markdown
        if mime in TEXT_MIMES:
            try:
                request = service.files().get_media(fileId=fid)
                buf = io.BytesIO()
                downloader = MediaIoBaseDownload(buf, request)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
                text = buf.getvalue().decode("utf-8", errors="ignore")
                return text, meta
            except Exception as e:
                print(f"Error downloading text file {name}: {e}")
                return "", meta

        # PDF: try pypdf first (fast), fallback to pdfminer
        if mime == PDF_MIME:
            try:
                request = service.files().get_media(fileId=fid)
                buf = io.BytesIO()
                downloader = MediaIoBaseDownload(buf, request)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
                data = io.BytesIO(buf.getvalue())
                try:
                    reader = PdfReader(data)
                    pages = [page.extract_text() or "" for page in reader.pages]
                    text = "\n\n".join(pages)
                except Exception:
                    data.seek(0)
                    text = pdf_extract(data)
                return text, meta
            except Exception as e:
                print(f"Error processing PDF {name}: {e}")
                return "", meta

        # Microsoft Office Documents (.docx, .xlsx, .pptx)
        if mime == DOCX_MIME or mime == DOC_MIME:
            try:
                request = service.files().get_media(fileId=fid)
                buf = io.BytesIO()
                downloader = MediaIoBaseDownload(buf, request)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
                
                if mime == DOCX_MIME:
                    text = extract_docx_text(buf.getvalue())
                else:
                    text = f"Legacy DOC file: {name} (text extraction not implemented)"
                
                return text, meta
            except Exception as e:
                print(f"Error processing Word document {name}: {e}")
                return "", meta

        if mime == XLSX_MIME or mime == XLS_MIME:
            try:
                request = service.files().get_media(fileId=fid)
                buf = io.BytesIO()
                downloader = MediaIoBaseDownload(buf, request)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
                
                if mime == XLSX_MIME:
                    text = extract_xlsx_text(buf.getvalue())
                else:
                    text = f"Legacy XLS file: {name} (text extraction not implemented)"
                
                return text, meta
            except Exception as e:
                print(f"Error processing Excel document {name}: {e}")
                return "", meta

        if mime == PPTX_MIME or mime == PPT_MIME:
            try:
                request = service.files().get_media(fileId=fid)
                buf = io.BytesIO()
                downloader = MediaIoBaseDownload(buf, request)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
                
                if mime == PPTX_MIME:
                    text = extract_pptx_text(buf.getvalue())
                else:
                    text = f"Legacy PPT file: {name} (text extraction not implemented)"
                
                return text, meta
            except Exception as e:
                print(f"Error processing PowerPoint document {name}: {e}")
                return "", meta

        # Unsupported types are skipped
        print(f"Unsupported file type: {mime} for file {name}")
        return "", meta
        
    except Exception as e:
        print(f"❌ Error processing file {file.get('name', 'Unknown')}: {e}")
        return "", {"id": file.get("id", ""), "name": file.get("name", "Unknown"), "mime": file.get("mimeType", "")}