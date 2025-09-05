# chat/gdrive_ingest.py — fetch & parse files from Google Drive using Service Account
import io
import json
from typing import Iterable, Tuple
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.oauth2 import service_account
from django.conf import settings
from pdfminer.high_level import extract_text as pdf_extract
from pypdf import PdfReader
from docx import Document
import zipfile

# Drive scopes: read-only
SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]

MIME_EXPORTS = {
    "application/vnd.google-apps.document": ("text/plain", "txt"),
    "application/vnd.google-apps.presentation": ("text/plain", "txt"),
    "application/vnd.google-apps.spreadsheet": ("text/plain", "txt"),
}

TEXT_MIMES = {
    "text/plain",
    "text/markdown",
}

PDF_MIME = "application/pdf"

# Microsoft Office document MIME types
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

# Legacy Office formats
DOC_MIME = "application/msword"
XLS_MIME = "application/vnd.ms-excel"
PPT_MIME = "application/vnd.ms-powerpoint"


def _creds():
    """Create credentials using Service Account for web app deployment."""
    credentials = service_account.Credentials.from_service_account_file(
        settings.GDRIVE_SERVICE_ACCOUNT_PATH,
        scopes=SCOPES
    )
    return credentials


def extract_docx_text(file_data: bytes) -> str:
    """Extract text from a .docx file."""
    try:
        doc = Document(io.BytesIO(file_data))
        paragraphs = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                paragraphs.append(paragraph.text.strip())
        return "\n\n".join(paragraphs)
    except Exception as e:
        print(f"Error extracting DOCX text: {e}")
        return ""


def extract_xlsx_text(file_data: bytes) -> str:
    """Extract text from an Excel file by reading all sheets and cells."""
    try:
        import openpyxl
        workbook = openpyxl.load_workbook(io.BytesIO(file_data), data_only=True)
        all_text = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = []
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                row_text = [cell for cell in row_text if cell.strip()]
                if row_text:
                    sheet_text.append(" | ".join(row_text))
            
            if sheet_text:
                all_text.append(f"Sheet: {sheet_name}\n" + "\n".join(sheet_text))
        
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"Error extracting XLSX text: {e}")
        return ""


def extract_pptx_text(file_data: bytes) -> str:
    """Extract text from a PowerPoint file."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_data))
        all_text = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                all_text.append(f"Slide {i}:\n" + "\n".join(slide_text))
        
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"Error extracting PPTX text: {e}")
        return ""



def iter_drive_files(folder_id: str) -> Iterable[dict]:
    """Iterate through all files in a Google Drive folder."""
    service = build("drive", "v3", credentials=_creds())
    q = f"'{folder_id}' in parents and trashed = false"
    fields = "nextPageToken, files(id, name, mimeType)"
    page_token = None
    while True:
        resp = service.files().list(q=q, fields=fields, pageToken=page_token).execute()
        for f in resp.get("files", []):
            yield f
        page_token = resp.get("nextPageToken")
        if not page_token:
            break


def fetch_text_for_file(file: dict) -> Tuple[str, dict]:
    """Returns (text, metadata) for a given Drive file."""
    service = build("drive", "v3", credentials=_creds())
    fid = file["id"]
    name = file["name"]
    mime = file["mimeType"]

    meta = {"id": fid, "name": name, "mime": mime}

    # Google Docs/Sheets/Slides: export as text
    if mime in MIME_EXPORTS:
        export_mime, ext = MIME_EXPORTS[mime]
        try:
            data = service.files().export(fileId=fid, mimeType=export_mime).execute()
            text = data.decode("utf-8", errors="ignore")
            return text, meta
        except Exception as e:
            print(f"Error exporting {name}: {e}")
            return "", meta

    # Plain text / Markdown
    if mime in TEXT_MIMES:
        try:
            request = service.files().get_media(fileId=fid)
            buf = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            text = buf.getvalue().decode("utf-8", errors="ignore")
            return text, meta
        except Exception as e:
            print(f"Error downloading text file {name}: {e}")
            return "", meta

    # PDF: try pypdf first (fast), fallback to pdfminer
    if mime == PDF_MIME:
        try:
            request = service.files().get_media(fileId=fid)
            buf = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            data = io.BytesIO(buf.getvalue())
            try:
                reader = PdfReader(data)
                pages = [page.extract_text() or "" for page in reader.pages]
                text = "\n\n".join(pages)
            except Exception:
                data.seek(0)
                text = pdf_extract(data)
            return text, meta
        except Exception as e:
            print(f"Error processing PDF {name}: {e}")
            return "", meta

    # Microsoft Office Documents (.docx, .xlsx, .pptx)
    if mime == DOCX_MIME or mime == DOC_MIME:
        try:
            request = service.files().get_media(fileId=fid)
            buf = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            
            if mime == DOCX_MIME:
                text = extract_docx_text(buf.getvalue())
            else:
                # For legacy .doc files, you might want to use python-docx2txt or similar
                text = f"Legacy DOC file: {name} (text extraction not implemented)"
            
            return text, meta
        except Exception as e:
            print(f"Error processing Word document {name}: {e}")
            return "", meta

    if mime == XLSX_MIME or mime == XLS_MIME:
        try:
            request = service.files().get_media(fileId=fid)
            buf = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            
            if mime == XLSX_MIME:
                text = extract_xlsx_text(buf.getvalue())
            else:
                text = f"Legacy XLS file: {name} (text extraction not implemented)"
            
            return text, meta
        except Exception as e:
            print(f"Error processing Excel document {name}: {e}")
            return "", meta

    if mime == PPTX_MIME or mime == PPT_MIME:
        try:
            request = service.files().get_media(fileId=fid)
            buf = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            
            if mime == PPTX_MIME:
                text = extract_pptx_text(buf.getvalue())
            else:
                text = f"Legacy PPT file: {name} (text extraction not implemented)"
            
            return text, meta
        except Exception as e:
            print(f"Error processing PowerPoint document {name}: {e}")
            return "", meta